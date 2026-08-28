"""Build slides/CS112_Team3_Presentation.pptx from slides/presentation.md.

A compact design system on top of python-pptx: a dark title slide and closing
slide, light content slides with a section kicker, an amber title rule, a
footer with page number, styled tables, tinted callout boxes and framed
charts. Body is laid out in document order with a running vertical cursor.

Run:  python slides/build_pptx.py       (needs:  pip install python-pptx)
"""

import math
import os
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))
MD = os.path.join(HERE, "presentation.md")
OUT = os.path.join(HERE, "CS112_Team3_Presentation.pptx")

# ---------------------------------------------------------------- design tokens
NAVY  = RGBColor(0x0E, 0x29, 0x42)
NAVY2 = RGBColor(0x17, 0x3A, 0x5B)
BLUE  = RGBColor(0x2F, 0x6F, 0xED)
TEAL  = RGBColor(0x15, 0xB8, 0xA6)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
INK   = RGBColor(0x1B, 0x2A, 0x38)
SLATE = RGBColor(0x5C, 0x6C, 0x7E)
MIST  = RGBColor(0xEC, 0xF2, 0xF8)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FOG   = RGBColor(0xB9, 0xCA, 0xDA)

FONT   = "Segoe UI"
FONT_SB = "Segoe UI Semibold"
FONT_L = "Segoe UI Light"

SW, SH = Inches(13.333), Inches(7.5)
ML = Inches(0.9)                     # left margin
CW_FULL = Inches(11.5)               # full content width
CW_TEXT = Inches(6.9)                # text col when a chart sits on the right
TEAM = "CS 112  ·  Cohort A — Team 3"
LINE_H = 0.34                        # inches per wrapped text line (approx)


# ---------------------------------------------------------------------- helpers
def _no_bullet(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for e in pPr.findall(qn(tag)):
            pPr.remove(e)
    pPr.append(etree.SubElement(pPr, qn("a:buNone")))


def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb, tf


def para(tf, first, align=PP_ALIGN.LEFT, spacing=1.12, after=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = spacing
    p.space_after = Pt(after)
    _no_bullet(p)
    return p


def run(p, s, size=16, color=INK, bold=False, italic=False, font=FONT):
    r = p.add_run()
    r.text = s
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return r


def wrapped_lines(s, chars_per_line):
    return max(1, math.ceil(len(s) / chars_per_line))


# ------------------------------------------------------------------- md parsing
def clean(t):
    t = re.sub(r"\*\*(.+?)\*\*", r"\1", t)
    t = re.sub(r"\*(.+?)\*", r"\1", t)
    t = re.sub(r"`(.+?)`", r"\1", t)
    return t.replace("→", "→").strip()


def parse(md):
    if md.startswith("---"):
        md = md.split("---", 2)[2]
    out = []
    for chunk in re.split(r"(?m)^---\s*$", md):
        chunk = re.sub(r"<!--.*?-->", "", chunk, flags=re.S).strip()
        if not chunk:
            continue
        lines = chunk.split("\n")
        title, level, body, i, tbl = None, 2, [], 0, []
        while i < len(lines):
            ln = lines[i].rstrip()
            h = re.match(r"^(#{1,6})\s+(.*)", ln)
            if h and title is None:
                level, title = len(h.group(1)), clean(h.group(2))
                i += 1
                continue
            img = re.match(r"^!\[[^\]]*\]\(([^)]+)\)", ln)
            if img:
                body.append(("image", img.group(1)))
                i += 1
                continue
            if ln.startswith("|"):
                tbl.append(ln)
                if i + 1 >= len(lines) or not lines[i + 1].rstrip().startswith("|"):
                    rows = []
                    for r in tbl:
                        cells = [c.strip() for c in r.strip().strip("|").split("|")]
                        if set("".join(cells)) <= set("-: "):
                            continue
                        rows.append([(c, "**" in c) for c in cells])
                    rows = [[(clean(c), b) for c, b in row] for row in rows]
                    if rows:
                        body.append(("table", rows))
                    tbl = []
                i += 1
                continue
            b = re.match(r"^(\s*)[-*]\s+(.*)", ln)
            if b:
                body.append(("bullet", clean(b.group(2)), len(b.group(1)) // 2))
                i += 1
                continue
            n = re.match(r"^\s*(\d+)\.\s+(.*)", ln)
            if n:
                body.append(("num", clean(n.group(2)), int(n.group(1))))
                i += 1
                continue
            raw = ln.strip()
            if raw:
                if re.match(r"^\*\*.+\*\*$", raw):
                    body.append(("callout", clean(raw)))
                elif re.match(r"^\*.+\*$", raw) or (raw.startswith("(") and raw.endswith(")")):
                    body.append(("note", clean(raw)))
                else:
                    body.append(("para", clean(raw)))
            i += 1
        out.append({"title": title, "level": level, "body": body})
    return out


def assign_kickers(slides):
    part_no, section = 0, "OVERVIEW"
    wrap_titles = {"testing", "what we learned", "live demo", "future work"}
    for s in slides:
        t = (s["title"] or "")
        m = re.match(r"^Part\s+(\d+)\s*[—-]\s*(.+)$", t, re.I)
        if m:
            part_no = int(m.group(1))
            rest = m.group(2).strip()
            name = "GRID ANALYSIS" if rest.lower() == "the dataset" else rest.upper()
            section = f"PART {part_no}  ·  {name}"
            s["title"] = rest
            s["kind"] = "content"
        elif s["level"] == 1:
            s["kind"] = "title"
        elif t.lower() == "thank you":
            s["kind"] = "closing"
        else:
            s["kind"] = "content"
            if t.lower() in wrap_titles:
                section = "WRAP-UP"
        s["kicker"] = None if s["kind"] != "content" else section
    return slides


# -------------------------------------------------------------------- renderers
def render_title(prs, s):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, SW, SH, NAVY)
    rect(sl, 0, 0, SW, Inches(0.18), TEAL)
    tb, tf = textbox(sl, ML, Inches(1.55), CW_FULL, Inches(0.4))
    run(para(tf, True), "CS 112  ·  COMPUTER PROGRAMMING FOR CS  ·  SUMMER 2026",
        size=12, color=FOG, bold=True)
    tb, tf = textbox(sl, ML, Inches(2.0), CW_FULL, Inches(2.0))
    p = para(tf, True, spacing=1.08, after=0)
    run(p, "National Grid Analysis", size=34, color=WHITE, bold=True, font=FONT_SB)
    p2 = para(tf, False, spacing=1.08)
    run(p2, "GridCare-Lite  ·  ClinicCare-Lite", size=34, color=TEAL, bold=True, font=FONT_SB)
    rect(sl, ML, Inches(4.25), Inches(2.4), Inches(0.09), TEAL)
    tb, tf = textbox(sl, ML, Inches(4.65), CW_FULL, Inches(2.4))
    skip = {"cs 112 — computer programming for cs · summer 2026"}
    first = True
    for it in s["body"]:
        t = it[1] if len(it) > 1 else ""
        if not t or t.lower() in skip:
            continue
        p = para(tf, first, spacing=1.3, after=8)
        first = False
        if t.startswith("github"):
            run(p, t, size=15, color=TEAL)
        elif "Team 3" in t:
            run(p, t, size=15, color=FOG)
        else:
            run(p, t, size=17, color=WHITE)
    return sl


def render_closing(prs, s):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, SW, SH, NAVY)
    rect(sl, 0, 0, SW, Inches(0.18), TEAL)
    tb, tf = textbox(sl, ML, Inches(2.7), CW_FULL, Inches(1.2))
    run(para(tf, True), s["title"], size=46, color=WHITE, bold=True, font=FONT_SB)
    rect(sl, ML, Inches(3.95), Inches(2.4), Inches(0.09), TEAL)
    tb, tf = textbox(sl, ML, Inches(4.35), CW_FULL, Inches(2))
    first = True
    for it in s["body"]:
        t = it[1] if len(it) > 1 else ""
        if not t:
            continue
        p = para(tf, first, spacing=1.35, after=6)
        first = False
        run(p, t, size=17, color=TEAL if t.startswith("github") else FOG)
    return sl


def render_table(sl, rows, x, y, w):
    nr = len(rows)
    nc = max(len(r) for r in rows)
    gt = sl.shapes.add_table(nr, nc, x, y, w, Inches(0.44 * nr)).table
    gt._tbl.find(qn("a:tblPr")).set("bandRow", "1")
    gt._tbl.find(qn("a:tblPr")).set("firstRow", "1")
    for ci in range(nc):
        gt.columns[ci].width = int(w / nc)
    for ri, row in enumerate(rows):
        gt.rows[ri].height = Inches(0.44)
        for ci in range(nc):
            c = gt.cell(ri, ci)
            c.margin_left = c.margin_right = Inches(0.14)
            c.margin_top = c.margin_bottom = Inches(0.04)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            txt, strong = row[ci] if ci < len(row) else ("", False)
            c.text = txt
            p = c.text_frame.paragraphs[0]
            _no_bullet(p)
            r0 = p.runs[0] if p.runs else p.add_run()
            r0.font.name = FONT
            if ri == 0:
                c.fill.solid(); c.fill.fore_color.rgb = NAVY
                r0.font.size = Pt(12); r0.font.bold = True; r0.font.color.rgb = WHITE
            else:
                c.fill.solid(); c.fill.fore_color.rgb = PAPER if ri % 2 else MIST
                r0.font.size = Pt(11.5); r0.font.color.rgb = INK; r0.font.bold = strong
    return 0.44 * nr


def render_content(prs, s, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, SW, SH, PAPER)
    rect(sl, 0, 0, Inches(0.16), SH, TEAL)
    if s["kicker"]:
        tb, tf = textbox(sl, ML, Inches(0.52), CW_FULL, Inches(0.3))
        run(para(tf, True), s["kicker"], size=11, color=TEAL, bold=True)
    tb, tf = textbox(sl, ML, Inches(0.86), CW_FULL, Inches(0.9))
    run(para(tf, True, spacing=1.0), s["title"], size=27, color=NAVY, bold=True, font=FONT_SB)
    rect(sl, ML, Inches(1.6), Inches(1.35), Inches(0.06), AMBER)
    # footer
    tb, tf = textbox(sl, ML, SH - Inches(0.52), Inches(8), Inches(0.3))
    run(para(tf, True), TEAM, size=9, color=SLATE)
    tb, tf = textbox(sl, SW - Inches(1.5), SH - Inches(0.52), Inches(0.9), Inches(0.3))
    run(para(tf, True, align=PP_ALIGN.RIGHT), str(page), size=9, color=SLATE, bold=True)

    body = s["body"]
    has_img = any(it[0] == "image" for it in body)
    has_tbl = any(it[0] == "table" for it in body)
    has_block = has_tbl or any(it[0] == "callout" for it in body)
    two_col = has_img and not has_tbl
    cw = CW_TEXT if two_col else CW_FULL
    cpl = 74 if two_col else 118      # chars per line estimate

    # pure-text slides: nudge down a little so the block sits optically centred
    y = 2.55 if not has_block else 2.05
    # group consecutive text items into single textboxes; tables/callouts/images inline
    pending = []

    def flush_text():
        nonlocal y
        if not pending:
            return
        h = sum(it["_h"] for it in pending) + 0.1
        tb, tf = textbox(sl, ML, Inches(y), cw, Inches(h))
        for k, it in enumerate(pending):
            p = para(tf, k == 0, spacing=1.14, after=it["_after"])
            if it["kind"] == "bullet":
                lvl = it["lvl"]
                run(p, ("▸  " if lvl == 0 else "–  "),
                    size=15 if lvl == 0 else 13, color=TEAL if lvl == 0 else FOG, bold=True)
                run(p, it["t"], size=15 if lvl == 0 else 13,
                    color=INK if lvl == 0 else SLATE)
                p.level = lvl
            elif it["kind"] == "num":
                run(p, f"{it['n']}.  ", size=15, color=BLUE, bold=True)
                run(p, it["t"], size=15, color=INK)
            elif it["kind"] == "note":
                run(p, it["t"], size=12, color=SLATE, italic=True)
            else:
                run(p, it["t"], size=14.5, color=INK)
        y += h + 0.12
        pending.clear()

    for it in body:
        kind = it[0]
        if kind in ("bullet", "num", "para", "note"):
            if kind == "bullet":
                d = {"kind": "bullet", "t": it[1], "lvl": it[2]}
                lines = wrapped_lines(("▸  " + it[1]), cpl if it[2] == 0 else cpl + 6)
            elif kind == "num":
                d = {"kind": "num", "n": it[2], "t": it[1]}
                lines = wrapped_lines(it[1], cpl - 3)
            else:
                d = {"kind": kind, "t": it[1]}
                lines = wrapped_lines(it[1], cpl)
            d["_h"] = lines * LINE_H
            d["_after"] = 4 if kind == "note" else 7
            pending.append(d)
        elif kind == "callout":
            flush_text()
            lines = wrapped_lines(it[1], 96)
            h = lines * 0.30 + 0.34
            rect(sl, ML, Inches(y), cw, Inches(h), MIST)
            rect(sl, ML, Inches(y), Inches(0.07), Inches(h), TEAL)
            tb, tf = textbox(sl, ML + Inches(0.28), Inches(y), cw - Inches(0.5), Inches(h),
                             anchor=MSO_ANCHOR.MIDDLE)
            run(para(tf, True, spacing=1.15), it[1], size=13.5, color=NAVY, bold=True)
            y += h + 0.18
        elif kind == "table":
            flush_text()
            dh = render_table(sl, it[1], ML, Inches(y), cw)
            y += dh + 0.2
        elif kind == "image":
            pass  # handled after, in right column
    flush_text()

    if has_img:
        img = next(it[1] for it in body if it[0] == "image")
        path = os.path.normpath(os.path.join(HERE, img))
        ix, iw = Inches(8.15), Inches(4.4)
        if os.path.exists(path):
            pic = sl.shapes.add_picture(path, ix, Inches(2.15), width=iw)
            fr = rect(sl, ix - Inches(0.05), Inches(2.15) - Inches(0.05),
                      iw + Inches(0.1), pic.height + Inches(0.1), MIST)
            sl.shapes._spTree.remove(fr._element)
            sl.shapes._spTree.insert(list(sl.shapes._spTree).index(pic._element), fr._element)
            tb, tf = textbox(sl, ix, Inches(2.15) + pic.height + Inches(0.12), iw, Inches(0.3))
            run(para(tf, True), os.path.basename(img), size=9, color=SLATE, italic=True)
        else:
            tb, tf = textbox(sl, ix, Inches(3.2), iw, Inches(1))
            run(para(tf, True), "[" + os.path.basename(img) + "]", size=11, color=SLATE)
    return sl


def main():
    with open(MD, encoding="utf-8") as f:
        slides = assign_kickers(parse(f.read()))
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    for i, s in enumerate(slides, 1):
        if s["kind"] == "title":
            render_title(prs, s)
        elif s["kind"] == "closing":
            render_closing(prs, s)
        else:
            render_content(prs, s, i)
    prs.save(OUT)
    print("wrote %s  (%d slides)" % (OUT, len(slides)))


if __name__ == "__main__":
    main()
